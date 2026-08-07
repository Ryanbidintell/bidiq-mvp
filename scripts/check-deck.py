#!/usr/bin/env python3
"""
scripts/check-deck.py — audit a Better Together .pptx against the constraints in
PROCORE_BETTER_TOGETHER_DECK.md before it goes to Procore.

These are the failure modes that are cheap to introduce by hand and expensive to have
found by a reviewer:
  * an invented customer quote (we have zero joint customers)
  * a bare win-rate percentage with no attribution
  * a "team analytics" claim (sold, not built)
  * write/bi-directional language, when the application answered "No" to bi-directional
    and the workflow diagram is read-only
  * Joint Solution bullets over Procore's 6-word cap
  * text blocks longer than Procore's "1-2 sentences" rule

Heuristic, not a proof. It reads text only — it cannot see layout, logo weight, or whether
a screenshot is legible. A clean run means "nothing obviously wrong in the copy," not
"ready to send."

Usage:  python scripts/check-deck.py path/to/deck.pptx
        python scripts/check-deck.py --self-test
"""

import re
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx not installed. Run: python -m pip install python-pptx")

# ── rules ────────────────────────────────────────────────────────────────────

ATTRIBUTION_TERMS = (
    'constructconnect', 'benchmark', 'illustrative', 'industry', 'source', 'target'
)
BANNED_PATTERNS = [
    (r'\bteam analytics\b',            "claims team analytics (sold but NOT built)"),
    (r'\bper-?estimator (analytics|metrics|reporting)\b',
                                       "claims per-estimator analytics (NOT built)"),
    (r'\bbi-?directional\b',           "says bi-directional; the application answered No"),
    (r'\bwrite(s|-)?\s?back\b(?!.*phase 2)', "mentions write-back without a phase-2 label"),
    (r'\b20\s?[-–]\s?25\s?%',          "bare 20-25% win-rate figure — unsourced, see deck 2a"),
]
QUOTE_RE = re.compile(r'["“”]([^"“”]{25,})["“”]')
SENTENCE_RE = re.compile(r'[.!?](?:\s|$)')
PERCENT_RE = re.compile(r'\b\d{1,3}\s?%')

SIX_WORD_SLIDE_HINT = 'joint solution'
MAX_BULLET_WORDS = 6
MAX_SENTENCES_PER_BLOCK = 2
MAX_BLOCK_CHARS = 240


def blocks(slide):
    """Yield (shape_name, text) for every non-empty text frame on a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            yield (shape.name, text)


def audit(path):
    prs = Presentation(path)
    findings = []

    def flag(slide_no, level, msg, excerpt=''):
        findings.append((slide_no, level, msg, excerpt[:110]))

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = ' '.join(t for _, t in blocks(slide)).lower()
        in_joint_solution = SIX_WORD_SLIDE_HINT in slide_text

        for _, text in blocks(slide):
            low = text.lower()

            for pattern, msg in BANNED_PATTERNS:
                if re.search(pattern, low):
                    flag(i, 'ERROR', msg, text)

            # Quoted strings long enough to be a testimonial.
            for q in QUOTE_RE.findall(text):
                if 'illustrative' not in low:
                    flag(i, 'ERROR',
                         "quoted passage with no 'illustrative' label — we have no joint customers",
                         q)

            # Percentages need an attribution word nearby.
            if PERCENT_RE.search(text) and not any(a in low for a in ATTRIBUTION_TERMS):
                flag(i, 'WARN', "percentage with no attribution or benchmark label", text)

            # Procore: 1-2 sentences per block.
            if len(SENTENCE_RE.findall(text)) > MAX_SENTENCES_PER_BLOCK:
                flag(i, 'WARN', "block runs past Procore's 1-2 sentence rule", text)
            if len(text) > MAX_BLOCK_CHARS:
                flag(i, 'WARN', f"block is {len(text)} chars — too long to read at a glance", text)

            # Joint Solution bullets are capped at 6 words.
            if in_joint_solution:
                for line in text.splitlines():
                    line = line.strip(' •-\t')
                    words = len(line.split())
                    if 0 < words <= 12 and words > MAX_BULLET_WORDS and line.endswith('.'):
                        flag(i, 'WARN',
                             f"possible Joint Solution bullet is {words} words (cap is {MAX_BULLET_WORDS})",
                             line)

    return prs, findings


def report(path):
    prs, findings = audit(path)
    print(f"\n{path} — {len(prs.slides)} slides\n" + "=" * 72)

    if not findings:
        print("No copy issues found.")
    else:
        for slide_no, level, msg, excerpt in sorted(findings, key=lambda f: (f[0], f[1])):
            mark = 'X' if level == 'ERROR' else '!'
            print(f"[{mark}] slide {slide_no}: {msg}")
            if excerpt:
                print(f"      {excerpt!r}")

    errors = sum(1 for f in findings if f[1] == 'ERROR')
    warns = len(findings) - errors
    print("=" * 72)
    print(f"{errors} error(s), {warns} warning(s)")
    print("\nStill check by hand — this reads text only:")
    print("  - logos at equal visual weight on slide 1")
    print("  - screenshots legible when projected (crop and zoom, don't scale down)")
    print("  - the pilot ask has a number attached")
    print("  - slides 1-3 carry the story on their own")
    return 1 if errors else 0


def self_test():
    """Build a deck that breaks the rules on purpose, and prove the checker catches it."""
    from pptx import Presentation as P
    import tempfile, os

    prs = P()
    blank = prs.slide_layouts[6]

    def add(text):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(0, 0, 5000000, 3000000)
        box.text_frame.text = text
        return s

    add('Better Together: BidIntell + Procore')
    add('"BidIntell cut our estimating hours in half in three weeks." — Acme Drywall')
    add('Subs win 20-25% of what they bid.')
    add('Includes team analytics for every estimator on your crew.')
    add('The integration is bi-directional and writes status back to Procore.')
    add('Joint Solution\nEstimating hours land on the work you can actually win.')
    add('ConstructConnect puts hard-bid work at 10-20%.')

    tmp = os.path.join(tempfile.gettempdir(), '_deck_selftest.pptx')
    prs.save(tmp)

    _, findings = audit(tmp)
    got = {(n, m) for n, lvl, m, _ in findings}

    expect = [
        (2, 'quote'), (3, '20-25'), (4, 'team analytics'),
        (5, 'bi-directional'), (6, 'Joint Solution bullet'),
    ]
    ok = True
    print("self-test")
    for slide_no, needle in expect:
        hit = any(n == slide_no and needle.lower() in m.lower() for n, m in got)
        print(f"  {'PASS' if hit else 'FAIL'}  slide {slide_no}: expected a finding matching {needle!r}")
        ok &= hit

    # Slide 7 attributes its percentage, so it must NOT be flagged for attribution.
    clean = not any(n == 7 and 'attribution' in m.lower() for n, m in got)
    print(f"  {'PASS' if clean else 'FAIL'}  slide 7: attributed percentage not flagged")
    ok &= clean

    os.remove(tmp)
    print("\nself-test", "PASSED" if ok else "FAILED")
    return 0 if ok else 1


if __name__ == '__main__':
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    if sys.argv[1] == '--self-test':
        sys.exit(self_test())
    sys.exit(report(sys.argv[1]))
